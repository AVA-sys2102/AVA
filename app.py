import streamlit as st
import json
import os
import io
from engine import AvaEngine
from pypdf import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation
from gtts import gTTS
from streamlit_mic_recorder import mic_recorder

# --- CONFIGURAÇÕES ---
st.set_page_config(page_title="AVA Advanced v4.5", page_icon="📑", layout="wide")
MEMORIA_FILE = "ava_memoria.json"

# --- PERSISTÊNCIA ---
def carregar_memoria():
    if os.path.exists(MEMORIA_FILE):
        try:
            with open(MEMORIA_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except: return []
    return []

def guardar_memoria(historico):
    with open(MEMORIA_FILE, "w", encoding="utf-8") as f:
        json.dump(historico, f, ensure_ascii=False, indent=4)

# --- INICIALIZAÇÃO ---
if "engine" not in st.session_state: st.session_state.engine = AvaEngine()
if "contexto_docs" not in st.session_state: st.session_state.contexto_docs = ""
if "historico" not in st.session_state: st.session_state.historico = carregar_memoria()

# --- MELHORIA NA LEITURA (PARSING ESTRUTURADO) ---
def extrair_texto_melhorado(files):
    texto_final = ""
    for f in files:
        ext = f.name.split(".")[-1].lower()
        texto_final += f"\n\n--- FICHEIRO: {f.name} ---\n" # Marcador para a IA
        
        try:
            if ext == "pdf":
                pdf = PdfReader(f)
                for i, page in enumerate(pdf.pages):
                    texto_final += f"[PÁGINA {i+1}]\n{page.extract_text()}\n"
            elif ext == "docx":
                doc = Document(f)
                texto_final += "\n".join([p.text for p in doc.paragraphs])
            elif ext == "xlsx":
                df = pd.read_excel(f)
                texto_final += f"Dados da Planilha:\n{df.to_string()}" # to_string é melhor que to_csv para leitura de IA
            elif ext == "pptx":
                prs = Presentation(f)
                for i, slide in enumerate(prs.slides):
                    texto_final += f"[SLIDE {i+1}]\n"
                    for shp in slide.shapes:
                        if hasattr(shp, "text"): texto_final += shp.text + "\n"
        except Exception as e:
            texto_final += f" Erro ao ler este ficheiro: {e}"
            
    return texto_final

def gerar_voz(texto):
    try:
        tts = gTTS(text=texto[:500], lang='pt') # Limitamos o áudio para ser mais rápido
        audio_fp = io.BytesIO()
        tts.write_to_fp(audio_fp)
        return audio_fp
    except: return None

# --- UI BARRA LATERAL ---
with st.sidebar:
    st.title("AVA - Advanced Virtual Assistent")
    
    st.subheader("1. Leitura Avançada")
    ficheiros = st.file_uploader("Carregar Office/PDF", accept_multiple_files=True)
    if st.button("📚 Analisar Documentos"):
        with st.spinner("A indexar conteúdos..."):
            st.session_state.contexto_docs = extrair_texto_melhorado(ficheiros)
            st.success("Documentos prontos para consulta!")

    st.divider()
    st.subheader("2. Visão")
    img_file = st.file_uploader("Ver Imagem", type=["jpg", "png", "jpeg"])
    if img_file: st.image(img_file, caption="Olhos da AVA")

    st.divider()
    st.subheader("3. Memória")
    memoria_on = st.toggle("Lembrar conversas anteriores", value=True)
    if st.button("🗑️ Reset Total"):
        if os.path.exists(MEMORIA_FILE): os.remove(MEMORIA_FILE)
        st.session_state.historico = []
        st.session_state.contexto_docs = ""
        st.rerun()

# --- ÁREA DE CHAT ---
st.title(" AVA - Advanced Virtual Assistent")

for msg in st.session_state.historico:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# Captura de Áudio e Texto
col1, col2 = st.columns([1, 4])
with col1:
    audio_data = mic_recorder(start_prompt="🎤 Falar", stop_prompt="⏹️ Enviar", key='mic')
with col2:
    prompt_txt = st.chat_input("Ou escreve a tua dúvida...")

# Processamento
if prompt_txt or audio_data:
    input_final = prompt_txt if prompt_txt else "Analisa o áudio que acabei de gravar."
    bytes_audio = audio_data['bytes'] if audio_data else None

    # Adicionar ao histórico
    st.session_state.historico.append({"role": "user", "content": input_final})
    with st.chat_message("user"): st.markdown(input_final)

    # Resposta
    with st.chat_message("assistant"):
        with st.spinner("A processar..."):
            bytes_img = img_file.getvalue() if img_file else None
            
            resposta = st.session_state.engine.processar_multimodal(
                mensagem_usuario=input_final,
                contexto_texto=st.session_state.contexto_docs,
                imagem_bytes=bytes_img,
                audio_bytes=bytes_audio,
                historico_anterior=st.session_state.historico if memoria_on else []
            )
            
            st.markdown(resposta)
            st.session_state.historico.append({"role": "assistant", "content": resposta})
            
            if memoria_on: guardar_memoria(st.session_state.historico)
            
            # Voz de Saída
            som = gerar_voz(resposta)
            if som: st.audio(som, format="audio/mp3", autoplay=True)